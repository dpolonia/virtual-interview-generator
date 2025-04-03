import os
import re
import json
import click
import questionary
import datetime
from rich.console import Console
from rich.table import Table
from rich.progress import Progress

from database.db_manager import db_manager
from models.ai_models import AIModelInterface, get_available_models, get_all_model_info, BatchProcessor
from prompts.prompt_templates import (
    INTERVIEWER_PERSONA_PROMPT,
    INTERVIEWEE_PERSONA_PROMPT,
    INTERVIEW_GENERATION_PROMPT,
    XML_FORMATTING_PROMPT,
    ANALYSIS_PROMPT
)
from utils.script_parser import (
    parse_interview_scripts,
    format_script_for_interview,
    get_category_key,
    save_scripts_to_files,
    STAKEHOLDER_CATEGORIES,
    NORMALIZED_CATEGORIES,
    save_scripts_to_json
)
from utils.persona_manager import FinePersonaManager

# Optional imports - will be handled with try/except when used
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.dml.color import RGBColor
# import markdown
# from weasyprint import HTML, CSS
# import pdfkit

console = Console()

@click.group()
def cli():
    """Virtual Interview Generator for AI in Consulting Research"""
    pass

def select_ai_model():
    """Interactive model selection with context window information."""
    # First, select provider
    provider_choices = [
        "OpenAI (GPT models)",
        "Anthropic (Claude models)",
        "Google (Gemini models)"
    ]
    
    provider = questionary.select(
        "Select AI provider:",
        choices=provider_choices
    ).ask()
    
    # Extract provider name
    provider_key = provider.split()[0].lower()
    
    # Get model information
    model_info = get_all_model_info()
    
    # Display model information
    table = Table(title=f"{provider} Models")
    table.add_column("Model Name")
    table.add_column("Description")
    table.add_column("Context Window")
    table.add_column("Output Tokens")
    table.add_column("Cost Tier")
    
    # Get only models for the selected provider
    provider_models = get_available_models()[provider_key]
    
    # Add rows for each model with context information
    for model_name in provider_models:
        info = model_info.get(model_name, {})
        context_size = format_token_count(info.get('context_window', 'Unknown'))
        output_tokens = format_token_count(info.get('token_limit', 'Unknown'))
        cost_tier = info.get('cost_tier', 'Unknown').capitalize()
        
        table.add_row(
            model_name, 
            info.get('description', 'No description'),
            context_size,
            output_tokens,
            cost_tier
        )
    
    console.print(table)
    
    # Add price guidance
    console.print("[yellow]Price Guidance:[/yellow]")
    console.print("• Budget: Lower cost, suitable for many interviews")
    console.print("• Standard: Mid-range pricing with good capabilities")
    console.print("• Premium: Higher cost with best quality/capabilities")
    
    # Add context window explanation
    console.print("\n[yellow]Context Window:[/yellow] Maximum tokens the model can process (input + output)")
    console.print("[yellow]Output Tokens:[/yellow] Maximum response length the model can generate")
    
    # Select specific model
    model_choices = get_available_models()[provider_key]
    model = questionary.select(
        "Select specific model:",
        choices=model_choices
    ).ask()
    
    return {"provider": provider_key, "model": model}

def format_token_count(count):
    """Format token count for display."""
    if isinstance(count, int):
        if count >= 1000:
            return f"{count/1000:.1f}K"
        return str(count)
    return str(count)

def create_presentation(presentation_bullets, output_dir, title, timestamp, filename):
    """
    Create a PowerPoint presentation from bullet points.
    
    Args:
        presentation_bullets: The formatted bullet points to include in the presentation
        output_dir: Directory to save the presentation
        title: Title for the presentation
        timestamp: Timestamp to include in the presentation
        filename: Filename for the presentation file (without extension)
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        console.print(f"[cyan]Creating PowerPoint presentation: {title}...[/cyan]")
        
        # Create a new presentation
        prs = Presentation()
        
        # Add a title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title_shape.text = title
        subtitle.text = f"Generated: {datetime.datetime.now().strftime('%Y-%m-%d')}"
        
        # Parse the bullet points into sections
        import re
        
        # Extract sections using regex
        sections = re.split(r'#+\s+', presentation_bullets)
        section_titles = re.findall(r'#+\s+(.*?)\n', "# " + presentation_bullets)
        
        # Create slides for each section
        for i, (section_content, section_title) in enumerate(zip(sections[1:], section_titles)):
            # Add a section slide
            bullet_slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            # Set the title
            slide_title = slide.shapes.title
            slide_title.text = section_title
            
            # Add the bullet points
            content = slide.placeholders[1]
            
            # Clean up the bullet points - remove empty lines and ensure proper formatting
            bullet_lines = []
            for line in section_content.strip().split('\n'):
                clean_line = line.strip()
                if clean_line and (clean_line.startswith('-') or clean_line.startswith('*')):
                    bullet_lines.append(clean_line[1:].strip())  # Remove the bullet character
            
            # Add to the content placeholder which automatically formats as bullets
            text_frame = content.text_frame
            text_frame.clear()  # Clear any default text
            
            for j, line in enumerate(bullet_lines):
                if j == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = line
                p.level = 0  # Top-level bullet
        
        # Add a thank you slide at the end
        thank_slide_layout = prs.slide_layouts[5]  # Title only layout
        slide = prs.slides.add_slide(thank_slide_layout)
        slide_title = slide.shapes.title
        slide_title.text = "Thank You"
        
        # Save the presentation
        pptx_path = os.path.join(output_dir, f'{filename}.pptx')
        prs.save(pptx_path)
        console.print(f"[green]PowerPoint presentation created at {pptx_path}[/green]")
        
        return pptx_path
        
    except ImportError:
        console.print("[yellow]Could not create PowerPoint presentation. Install 'python-pptx' package if needed.[/yellow]")
        console.print("Run: pip install python-pptx")
        
        # Create a simplified version as markdown that can be easily converted
        ppt_md_path = os.path.join(output_dir, f'{filename}.md')
        with open(ppt_md_path, 'w', encoding='utf-8') as f:
            f.write(f"# {title}\n\n")
            f.write(f"Generated: {datetime.datetime.now().strftime('%Y-%m-%d')}\n\n")
            f.write(presentation_bullets)
            f.write("\n\n# Thank You")
        
        console.print(f"[green]Created presentation markdown at {ppt_md_path}[/green]")
        
        return ppt_md_path

def generate_executive_summary(ai_model, analyses, category=None):
    """
    Generate an executive summary from multiple interview analyses.
    
    Args:
        ai_model: The AI model to use for generation
        analyses: List of analysis objects to summarize
        category: Optional category name for category-specific summaries
    
    Returns:
        dict containing the executive_summary and presentation_bullets
    """
    if not analyses:
        return {
            "executive_summary": "No analyses available to generate an executive summary.",
            "presentation_bullets": "No analyses available to generate presentation points."
        }
    
    # Extract key points and insights by research question
    key_points = []
    rq_insights = {
        "rq1": [],
        "rq2": [],
        "rq3": [],
        "rq4": []
    }
    
    for analysis in analyses:
        if analysis.key_points:
            key_points.append(analysis.key_points)
        
        if analysis.rq1_insights:
            rq_insights["rq1"].append(analysis.rq1_insights)
        if analysis.rq2_insights:
            rq_insights["rq2"].append(analysis.rq2_insights)
        if analysis.rq3_insights:
            rq_insights["rq3"].append(analysis.rq3_insights)
        if analysis.rq4_insights:
            rq_insights["rq4"].append(analysis.rq4_insights)
    
    # Build the executive summary prompt
    category_text = f"for {category.replace('_', ' ')} stakeholders" if category else "across all stakeholder perspectives"
    
    exec_summary_prompt = f"""
Generate an executive summary for a research report on "The Role of Business Consulting Firms in the Era of Artificial Intelligence" 
based on the following key points extracted from {len(analyses)} interviews {category_text}.

KEY POINTS FROM INTERVIEWS:
{"".join(key_points[:5])}  # Use first 5 to avoid token limits

Your executive summary should:
1. Synthesize the main findings {category_text}
2. Highlight the most significant trends and insights
3. Discuss implications for the consulting industry
4. Be approximately 500 words in length
"""
    
    # Generate the executive summary
    executive_summary = ai_model.generate_text(exec_summary_prompt)
    
    # Generate presentation-ready bullet points
    presentation_prompt = f"""
Based on the executive summary below and research insights from {len(analyses)} interviews about AI in consulting {category_text},
create concise, presentation-ready bullet points organized by key sections.

EXECUTIVE SUMMARY:
{executive_summary}

For each of the following categories, create 3-5 crisp, insightful bullet points suitable for a presentation slide:

1. Overall Key Findings
2. AI Adoption Status (RQ1)
3. Market Trends (RQ2)
4. Automation & Knowledge Effects (RQ3)
5. Ethical Considerations (RQ4)
"""
    
    # Add category-specific sections as appropriate
    if category:
        presentation_prompt += f"6. Key Implications for {category.replace('_', ' ').title()}\n"
    else:
        presentation_prompt += "6. Recommendations for Consulting Firms\n"
    
    presentation_prompt += """
Each bullet should be:
- Specific and evidence-based
- 1-2 lines maximum
- Clear and impactful for a business audience
- Free of jargon and unnecessary qualifiers
- Formatted for direct placement in presentation slides

Format with clear headers and bullet points using * or - symbols.
"""
    
    presentation_bullets = ai_model.generate_text(presentation_prompt)
    
    return {
        "executive_summary": executive_summary,
        "presentation_bullets": presentation_bullets
    }

@cli.command()
def preload_personas():
    """Preload a local cache of FinePersonas for faster searching."""
    console.print("[bold]FinePersonas Cache Generator[/bold]")
    console.print("This will preload a subset of personas from FinePersonas dataset for faster searching.")
    
    # Ask for dataset source
    use_sample = questionary.select(
        "Which dataset would you like to use?",
        choices=[
            "Sample dataset (100k personas - faster but less diverse)",
            "Full dataset (21M personas - much slower but more comprehensive)"
        ]
    ).ask().startswith("Sample")
    
    # Initialize the persona manager
    persona_manager = FinePersonaManager(use_sample=use_sample)
    
    # Ask if we should enable semantic search for higher quality matches
    use_semantic = questionary.confirm(
        "Enable semantic search for higher quality matches? (Requires scikit-learn)",
        default=True
    ).ask()
    
    # Ask for personas per category
    console.print("\n[bold]Number of personas to preload for each category:[/bold]")
    
    category_counts = {}
    
    # Get count for stakeholder categories
    for category_name, category_key in NORMALIZED_CATEGORIES.items():
        count = questionary.text(
            f"How many personas for {category_name}?",
            default="5000"
        ).ask()
        try:
            count = int(count)
            category_counts[category_key] = count
        except ValueError:
            console.print(f"[yellow]Invalid number, using default of 5000 for {category_name}[/yellow]")
            category_counts[category_key] = 5000
    
    # Get count for interviewers
    count = questionary.text(
        "How many personas for interviewers?",
        default="5000"
    ).ask()
    try:
        count = int(count)
        category_counts["interviewer"] = count
    except ValueError:
        console.print("[yellow]Invalid number, using default of 5000 for interviewers[/yellow]")
        category_counts["interviewer"] = 5000
    
    # Calculate total and ask for confirmation
    total_personas = sum(category_counts.values())
    console.print(f"\nWill preload a total of {total_personas} personas.")
    
    proceed = questionary.confirm(
        "Proceed with preloading? This may take some time.",
        default=True
    ).ask()
    
    if proceed:
        # Load the dataset
        if persona_manager.dataset is None:
            console.print("[cyan]Loading FinePersonas dataset...[/cyan]")
            persona_manager.load_dataset()
        
        # Preload personas
        console.print("[cyan]Preloading personas by category...[/cyan]")
        # Initialize semantic search if requested
        if use_semantic:
            persona_manager.initialize_vectorizer()
            
        # Preload personas
        success = persona_manager.preload_personas_by_categories(category_counts, max_total=total_personas)
        
        if success:
            console.print("[green]Successfully preloaded personas! Future searches will be much faster.[/green]")
        else:
            console.print("[red]Error preloading personas.[/red]")
    else:
        console.print("Preloading cancelled.")

@cli.command()
def parse_scripts():
    """Manage JSON interview scripts."""
    # Ask user if they want to use an existing JSON file or create/edit one
    source_type = questionary.select(
        "Select action:",
        choices=["Use existing JSON file", "Create/edit JSON file"]
    ).ask()
    
    if source_type == "Create/edit JSON file":
        # This option allows manually creating or editing the JSON file structure
        console.print("This will guide you through creating or editing the interview questions JSON file.")
        
        json_path = "data/scripts/interview_questions.json"
        data = {}
        
        # Check if file already exists
        if os.path.exists(json_path):
            try:
                with open(json_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                console.print("[green]Loaded existing questions file.[/green]")
            except Exception as e:
                console.print(f"[yellow]Error loading existing file: {str(e)}. Starting fresh.[/yellow]")
        
        # Get list of categories to work with
        categories_to_edit = questionary.checkbox(
            "Select categories to create or edit:",
            choices=[cat for cat in NORMALIZED_CATEGORIES.values()]
        ).ask()
        
        # Edit each selected category
        for category_key in categories_to_edit:
            full_name = next((cat for cat in NORMALIZED_CATEGORIES.keys() 
                           if NORMALIZED_CATEGORIES[cat] == category_key), category_key)
            
            console.print(f"\n[bold]Editing questions for {full_name}[/bold]")
            
            # Initialize category if it doesn't exist
            if category_key not in data:
                data[category_key] = {
                    "name": full_name,
                    "questions": {
                        "Demographic": [],
                        "RQ1": [],
                        "RQ2": [],
                        "RQ3": [],
                        "RQ4": []
                    }
                }
            
            # Edit questions for each section
            for section in ["Demographic", "RQ1", "RQ2", "RQ3", "RQ4"]:
                console.print(f"\n[cyan]{section} Questions:[/cyan]")
                
                # Show existing questions
                if section in data[category_key]["questions"] and data[category_key]["questions"][section]:
                    for i, q in enumerate(data[category_key]["questions"][section]):
                        console.print(f"{i+1}. {q}")
                else:
                    console.print("[yellow]No questions defined yet.[/yellow]")
                
                # Ask if user wants to edit this section
                if questionary.confirm(f"Edit {section} questions?").ask():
                    # Create a new list for this section
                    new_questions = []
                    
                    # Option to add questions one by one
                    adding = True
                    while adding:
                        new_q = questionary.text("Enter question (leave empty to finish):").ask()
                        if not new_q:
                            adding = False
                        else:
                            new_questions.append(new_q)
                    
                    # Update the questions for this section
                    data[category_key]["questions"][section] = new_questions
        
        # Save the updated data
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=2)
        
        console.print("[green]Questions saved successfully to data/scripts/interview_questions.json[/green]")
    
    else:
        # Use existing JSON file
        json_path = questionary.text(
            "Enter the path to the JSON file:",
            default="data/scripts/interview_questions.json"
        ).ask()
        
        console.print(f"Using scripts from {json_path}...")
        
        try:
            # Verify JSON file exists and is valid
            with open(json_path, 'r', encoding='utf-8') as f:
                scripts_data = json.load(f)
            
            # Display categories found in JSON
            table = Table(title="Available Interview Scripts")
            table.add_column("Category")
            table.add_column("Questions")
            
            for category_key, data in scripts_data.items():
                category_name = data.get("name", category_key)
                question_count = sum(len(qs) for qs in data.get("questions", {}).values()) if "questions" in data else 0
                table.add_row(category_name, str(question_count))
            
            console.print(table)
            console.print("[green]Scripts are ready to use.[/green]")
            
        except Exception as e:
            console.print(f"[red]Error loading JSON file: {str(e)}[/red]")

@cli.command()
def generate_personas():
    """Generate interviewer and interviewee personas."""
    # Ask if user wants to use FinePersonas
    use_finepersonas = questionary.confirm(
        "Do you want to use FinePersonas dataset for persona generation?",
        default=True
    ).ask()
    
    if use_finepersonas:
        # Initialize FinePersona manager
        persona_manager = FinePersonaManager(use_sample=True)  # Use the smaller sample by default
        
        # Load the dataset
        if not persona_manager.load_dataset():
            console.print("[yellow]Failed to load FinePersonas dataset. Falling back to AI generation.[/yellow]")
            use_finepersonas = False
    
    if use_finepersonas:
        # Generate interviewer personas from FinePersonas
        num_interviewers = questionary.text(
            "How many interviewer personas do you want to generate?",
            default="5"
        ).ask()
        num_interviewers = int(num_interviewers)
        
        # Use 'industry_analysts' as a good source for interviewers
        # Ask if the user wants to search for specific interviewer characteristics
        use_search = questionary.confirm(
            "Do you want to search for specific interviewer characteristics?",
            default=False
        ).ask()
        
        if use_search:
            # Ask whether to use the full dataset (better results but slower)
            use_full_dataset = questionary.confirm(
                "Use the full FinePersonas dataset instead of the sample? (Better results but slower and requires 143GB)",
                default=False
            ).ask()
            
            search_query = questionary.text(
                "Enter search terms for interviewer (e.g., 'experienced researcher AI ethics'):"
            ).ask()
            
            # Create or update the persona manager with the chosen dataset size
            if not 'persona_manager' in locals():
                persona_manager = FinePersonaManager(use_sample=not use_full_dataset)
            else:
                persona_manager.use_sample = not use_full_dataset
                persona_manager.dataset_size = "100k sample" if not use_full_dataset else "full 21M"
                
            # Load dataset if needed
            if persona_manager.dataset is None:
                console.print("[cyan]Loading FinePersonas dataset...[/cyan]")
                persona_manager.load_dataset()
                
            interviewer_personas = persona_manager.get_personas_by_category('industry_analysts', num_interviewers, search_query)
        else:
            interviewer_personas = persona_manager.get_personas_by_category('industry_analysts', num_interviewers)
        
        console.print(f"Selected {len(interviewer_personas)} interviewer personas from FinePersonas")
        
        # Save interviewers to database
        for persona_data in interviewer_personas:
            formatted_persona = persona_manager.format_persona_for_interview(persona_data, 'interviewer')
            db_manager.create_persona(formatted_persona)
        
        # Ask which categories to generate interviewees for
        categories = []
        while not categories:
            categories = questionary.checkbox(
                "Select at least one stakeholder category to generate personas for (IMPORTANT: You must select at least one):",
                choices=[
                    'senior_executives',
                    'ai_specialists',
                    'mid_level_consultants',
                    'clients',
                    'technology_providers',
                    'regulatory_stakeholders',
                    'industry_analysts'
                ]
            ).ask()
            
            if not categories:
                console.print("[red]ERROR: Please select at least one category. You did not select any categories.[/red]")
        
        # Ask how many personas per category
        num_per_category = questionary.text(
            "How many personas to generate per category?",
            default="10"
        ).ask()
        num_per_category = int(num_per_category)
        
        # Generate interviewee personas for each category
        for category in categories:
            console.print(f"\nGenerating personas for {category}...")
            
            # Ask if the user wants to search for specific characteristics for this category
            use_category_search = questionary.confirm(
                f"Do you want to search for specific characteristics in {category} personas?",
                default=False
            ).ask()
            
            if use_category_search:
                search_query = questionary.text(
                    f"Enter search terms for {category} personas:"
                ).ask()
                category_personas = persona_manager.get_personas_by_category(category, num_per_category, search_query)
            else:
                category_personas = persona_manager.get_personas_by_category(category, num_per_category)
            console.print(f"Selected {len(category_personas)} personas from FinePersonas")
            
            # Save personas to database
            for persona_data in category_personas:
                formatted_persona = persona_manager.format_persona_for_interview(persona_data, 'interviewee')
                formatted_persona['category'] = category  # Ensure correct category
                db_manager.create_persona(formatted_persona)
        
        console.print("[green]All personas generated successfully from FinePersonas![/green]")
        
    else:
        # Select AI model for persona generation
        console.print("Selecting AI model for persona generation:")
        model_info = select_ai_model()
        
        try:
            ai_model = AIModelInterface(model_info["provider"], model_info["model"])
            
            # Ask how many interviewer personas to generate
            num_interviewers = questionary.text(
                "How many interviewer personas do you want to generate?",
                default="5"
            ).ask()
            num_interviewers = int(num_interviewers)
            
            # Generate interviewer personas
            console.print(f"Generating {num_interviewers} interviewer personas...")
            
            with Progress() as progress:
                task = progress.add_task("[green]Generating interviewer personas...", total=num_interviewers)
                
                for i in range(num_interviewers):
                    response = ai_model.generate_text(INTERVIEWER_PERSONA_PROMPT)
                    
                    # Parse the response to extract persona details
                    # This is a simplified version - you might want to add more structured parsing
                    persona_data = {
                        "name": f"Interviewer {i+1}",  # Placeholder - extract from response
                        "category": "interviewer",
                        "role": "interviewer",
                        "background": response,
                        "created_by": f"{model_info['provider']}/{model_info['model']}"
                    }
                    
                    # Save to database
                    db_manager.create_persona(persona_data)
                    progress.update(task, advance=1)
            
            # Ask which categories to generate interviewees for
            categories = []
            while not categories:
                categories = questionary.checkbox(
                    "Select at least one stakeholder category to generate personas for (IMPORTANT: You must select at least one):",
                    choices=STAKEHOLDER_CATEGORIES
                ).ask()
                
                if not categories:
                    console.print("[red]ERROR: Please select at least one category. You did not select any categories.[/red]")
            
            # Ask how many personas per category
            num_per_category = questionary.text(
                "How many personas to generate per category?",
                default="10"
            ).ask()
            num_per_category = int(num_per_category)
            
            # Generate interviewee personas
            for category in categories:
                console.print(f"\nGenerating {num_per_category} personas for {category}...")
                
                with Progress() as progress:
                    task = progress.add_task(f"[green]Generating {category} personas...", total=num_per_category)
                    
                    for i in range(num_per_category):
                        prompt = INTERVIEWEE_PERSONA_PROMPT.format(category=category)
                        response = ai_model.generate_text(prompt)
                        
                        # Parse the response to extract persona details
                        # This is a simplified version - you might want to add more structured parsing
                        persona_data = {
                            "name": f"{category} {i+1}",  # Placeholder - extract from response
                            "category": get_category_key(category),
                            "role": "interviewee",
                            "background": response,
                            "created_by": f"{model_info['provider']}/{model_info['model']}"
                        }
                        
                        # Save to database
                        db_manager.create_persona(persona_data)
                        progress.update(task, advance=1)
            
            console.print("[green]All personas generated successfully![/green]")
            
        except Exception as e:
            console.print(f"[red]Error generating personas: {str(e)}[/red]")

@cli.command()
def generate_interviews():
    """Generate interviews between selected personas across all stakeholder categories."""
    # Track which categories have been interviewed
    interviewed_categories = set()
    
    # Track total interviews generated
    total_interviews_generated = 0
    
    # Continue interviewing until user decides to stop
    continue_interviewing = True
    while continue_interviewing:
        # Get remaining categories that haven't been interviewed yet in this round
        remaining_categories = [cat for cat in STAKEHOLDER_CATEGORIES 
                               if get_category_key(cat) not in interviewed_categories]
        
        # If all categories have been interviewed in this round, ask if user wants another round
        if not remaining_categories:
            console.print(f"[green]All {len(STAKEHOLDER_CATEGORIES)} stakeholder categories have been interviewed![/green]")
            
            # Reset for a new round if user wants to continue
            another_round = questionary.confirm(
                "Would you like to conduct another round of interviews?",
                default=True
            ).ask()
            
            if another_round:
                interviewed_categories.clear()
                remaining_categories = STAKEHOLDER_CATEGORIES
                console.print("[cyan]Starting a new round of interviews...[/cyan]")
            else:
                console.print(f"[green]Interview generation complete. Total interviews: {total_interviews_generated}[/green]")
                
                # Ask if user wants to generate a report
                generate_report_now = questionary.confirm(
                    "Would you like to generate a comprehensive report of all interviews now?",
                    default=True
                ).ask()
                
                if generate_report_now:
                    console.print("[cyan]Generating comprehensive report...[/cyan]")
                    # Call the report generation function
                    generate_report()
                
                # Exit the interview loop
                break
        
        # Select AI model for interview generation (only once per round)
        if not 'model_info' in locals() or not model_info:
            console.print("Selecting AI model for interview generation:")
            model_info = select_ai_model()
    
    try:
        ai_model = AIModelInterface(model_info["provider"], model_info["model"])
        
        # Get stakeholder categories to interview in this session
        if len(remaining_categories) > 1:
            # Allow selecting multiple categories at once
            selected_categories = questionary.checkbox(
                "Select stakeholder categories to interview in this session:",
                choices=remaining_categories
            ).ask()
            
            if not selected_categories:
                console.print("[yellow]No categories selected. Please select at least one category.[/yellow]")
                continue
        else:
            # If only one category remains, use it directly
            selected_categories = remaining_categories
        
        # Process each selected category
        for category_name in selected_categories:
            console.print(f"\n[bold cyan]Generating interviews for: {category_name}[/bold cyan]")
            category_interviews_generated = 0
            category_key = get_category_key(category_name)
        
        # Load script for this category from the JSON file
        json_path = "data/scripts/interview_questions.json"
        if not os.path.exists(json_path):
            console.print(f"[yellow]Script file not found. Run 'parse_scripts' command first.[/yellow]")
            return
        
        try:
            with open(json_path, 'r', encoding='utf-8') as f:
                scripts_data = json.load(f)
            
            if category_key not in scripts_data:
                console.print(f"[yellow]Category {category_name} not found in scripts file.[/yellow]")
                return
                
            # Format the script from the JSON structure
            script_text = ""
            if "questions" in scripts_data[category_key]:
                for section, questions in scripts_data[category_key]["questions"].items():
                    script_text += f"{section} Questions:\n"
                    for question in questions:
                        script_text += f"- {question}\n"
                    script_text += "\n"
            else:
                script_text = scripts_data[category_key].get("script", "")
                
            script = format_script_for_interview(script_text)
        except Exception as e:
            console.print(f"[red]Error loading script: {str(e)}[/red]")
            return
        
        # Get interviewer personas
        interviewers = db_manager.get_personas_by_category("interviewer", "interviewer")
        if not interviewers:
            console.print("[yellow]No interviewer personas found. Run 'generate_personas' first.[/yellow]")
            return
        
        # Create more descriptive choices to better represent personas
        interviewer_choices = []
        for p in interviewers:
            # Extract a short description from the background (first sentence or part of it)
            short_desc = ""
            if p.background:
                first_sentence = p.background.split('.')[0]
                short_desc = (first_sentence[:100] + '...') if len(first_sentence) > 100 else first_sentence
            interviewer_choices.append(f"{p.name} - {short_desc}")
        
        # Check if we should search FinePersonas for more matching personas
        use_finepersonas_search = questionary.confirm(
            "Do you want to search for additional interviewer personas in FinePersonas?",
            default=False
        ).ask()
        
        if use_finepersonas_search:
            # Generate suggested queries for interviewers based on the selected category
            suggested_queries = [
                "researcher AI consulting",
                "professor business technology",
                "academic specializing in artificial intelligence",
                "industry analyst mckinsey bcg",
                "technology journalist consulting",
                "consultant gartner forrester analyst",
                "senior advisor deloitte accenture"
            ]
            
            # Add category-specific suggestions
            if category_key == 'senior_executives':
                suggested_queries.append("researcher executive leadership mckinsey bcg")
            elif category_key == 'ai_specialists':
                suggested_queries.append("researcher artificial intelligence accenture palantir")
            elif category_key == 'regulatory_stakeholders':
                suggested_queries.append("professor policy regulation kpmg deloitte")
            
            # Allow user to select from suggested queries or create their own
            query_choices = suggested_queries + ["Enter my own search terms"]
            
            query_selection = questionary.select(
                "Select a suggested search query or create your own:",
                choices=query_choices
            ).ask()
            
            # Get the final search query
            if query_selection == "Enter my own search terms":
                search_query = questionary.text(
                    "Enter search terms for interviewer:"
                ).ask()
            else:
                search_query = query_selection
                console.print(f"Using selected query: '{search_query}'")
            
            console.print("Searching FinePersonas for matching interviewer profiles...")
            
            # Initialize FinePersona manager if not already done
            if not 'persona_manager' in locals():
                persona_manager = FinePersonaManager()
                console.print("[cyan]Loading FinePersonas cached personas...[/cyan]")
                persona_manager.load_dataset()
                
                # If no cache exists, ask if user wants to create one
                if not any(persona_manager.category_personas.values()):
                    create_cache = questionary.confirm(
                        "No cached personas found. Would you like to preload personas now? (Recommended for faster searches)",
                        default=True
                    ).ask()
                    
                    if create_cache:
                        console.print("[cyan]Running preload command...[/cyan]")
                        # Call the preload command
                        preload_personas()
            
            # Get matching personas
            search_results = persona_manager.get_personas_by_category('industry_analysts', 5, search_query)
            
            if search_results:
                # Format the personas and add to database
                for persona_data in search_results:
                    formatted_persona = persona_manager.format_persona_for_interview(persona_data, 'interviewer')
                    db_manager.create_persona(formatted_persona)
                
                # Refresh interviewer list
                interviewers = db_manager.get_personas_by_category("interviewer", "interviewer")
                interviewer_choices = []
                for p in interviewers:
                    short_desc = ""
                    if p.background:
                        first_sentence = p.background.split('.')[0]
                        short_desc = (first_sentence[:100] + '...') if len(first_sentence) > 100 else first_sentence
                    interviewer_choices.append(f"{p.name} - {short_desc}")
                
                console.print(f"[green]Added {len(search_results)} new interviewer profiles.[/green]")
        
        selected_interviewer_display = questionary.select(
            "Select interviewer persona:",
            choices=interviewer_choices
        ).ask()
        
        # Extract just the name part
        selected_interviewer_name = selected_interviewer_display.split(' - ')[0]
        selected_interviewer = next(p for p in interviewers if p.name == selected_interviewer_name)
        
        # Get interviewee personas for this category
        interviewees = db_manager.get_personas_by_category(category_key, "interviewee")
        if not interviewees:
            console.print(f"[yellow]No personas found for category {category_name}. Run 'generate_personas' first and select this category.[/yellow]")
            return
        
        # Create more descriptive choices for interviewees too
        interviewee_choices = []
        for p in interviewees:
            # Extract a short description from the background (first sentence or part of it)
            short_desc = ""
            if p.background:
                first_sentence = p.background.split('.')[0]
                short_desc = (first_sentence[:100] + '...') if len(first_sentence) > 100 else first_sentence
            interviewee_choices.append(f"{p.name} - {short_desc}")
        
        # Check if we should search FinePersonas for more matching interviewee personas
        use_finepersonas_search = questionary.confirm(
            f"Do you want to search for additional {category_name} personas in FinePersonas?",
            default=False
        ).ask()
        
        if use_finepersonas_search:
            # Generate suggested query based on the category
            suggested_query = ""
            
            # Load the JSON file to get category-specific details
            try:
                with open("data/scripts/interview_questions.json", 'r', encoding='utf-8') as f:
                    scripts_data = json.load(f)
                
                if category_key in scripts_data and "questions" in scripts_data[category_key]:
                    # Extract key terms from questions for this category
                    all_questions = []
                    for section, questions in scripts_data[category_key]["questions"].items():
                        all_questions.extend(questions)
                    
                    # Find common terms in questions that could be relevant for search
                    import re
                    from collections import Counter
                    
                    # Extract key terms (nouns and adjectives likely to be profile indicators)
                    words = re.findall(r'\b([A-Z][a-z]{2,}|[a-z]{3,})\b', ' '.join(all_questions))
                    word_count = Counter(words)
                    
                    # Filter out common words that aren't useful for profile search
                    common_words = {'have', 'what', 'how', 'your', 'with', 'consulting', 'about', 'that', 'this',
                                   'industry', 'technologies', 'specific', 'most', 'think', 'would', 'firm'}
                    
                    # Build a suggested query from the most frequent relevant terms
                    key_terms = [word for word, count in word_count.most_common(10) 
                                if count > 1 and word.lower() not in common_words][:3]
                    
                    # Create suggestion based on category with consulting firm integration
                    if category_key == 'senior_executives':
                        # For senior executives, add top consulting firms
                        suggested_query = f"executive consulting mckinsey bcg bain {' '.join(key_terms)}"
                    elif category_key == 'ai_specialists':
                        # For AI specialists, add tech-focused firms
                        suggested_query = f"AI specialist consulting accenture palantir {' '.join(key_terms)}"
                    elif category_key == 'mid_level_consultants':
                        # For mid-level, include a mix of firms
                        suggested_query = f"consultant deloitte mckinsey accenture {' '.join(key_terms)}"
                    elif category_key == 'clients':
                        suggested_query = f"business client {' '.join(key_terms)}"
                    elif category_key == 'technology_providers':
                        suggested_query = f"technology provider accenture ibm {' '.join(key_terms)}"
                    elif category_key == 'regulatory_stakeholders':
                        suggested_query = f"regulation deloitte ey kpmg {' '.join(key_terms)}"
                    elif category_key == 'industry_analysts':
                        suggested_query = f"analyst gartner forrester {' '.join(key_terms)}"
            except Exception as e:
                console.print(f"[yellow]Error generating suggested query: {str(e)}[/yellow]")
                suggested_query = f"{category_name.lower().replace('_', ' ')}"
            
            # Offer the suggested query to the user
            use_suggested = False
            if suggested_query:
                use_suggested = questionary.confirm(
                    f"Use suggested search query: '{suggested_query}'?",
                    default=True
                ).ask()
            
            # Get the final search query
            if use_suggested and suggested_query:
                search_query = suggested_query
                console.print(f"Using suggested query: '{search_query}'")
            else:
                search_query = questionary.text(
                    f"Enter search terms for {category_name} personas:",
                    default=suggested_query
                ).ask()
            
            console.print(f"Searching FinePersonas for matching {category_name} profiles...")
            
            # Initialize FinePersona manager if not already done
            if not 'persona_manager' in locals():
                persona_manager = FinePersonaManager()
                console.print("[cyan]Loading FinePersonas cached personas...[/cyan]")
                persona_manager.load_dataset()
                
                # If no cache exists, ask if user wants to create one
                if not any(persona_manager.category_personas.values()):
                    create_cache = questionary.confirm(
                        "No cached personas found. Would you like to preload personas now? (Recommended for faster searches)",
                        default=True
                    ).ask()
                    
                    if create_cache:
                        console.print("[cyan]Running preload command...[/cyan]")
                        # Call the preload command
                        preload_personas()
            
            # Get matching personas
            search_results = persona_manager.get_personas_by_category(category_key, 5, search_query)
            
            if search_results:
                # Format the personas and add to database
                for persona_data in search_results:
                    formatted_persona = persona_manager.format_persona_for_interview(persona_data, 'interviewee')
                    formatted_persona['category'] = category_key  # Ensure correct category
                    db_manager.create_persona(formatted_persona)
                
                # Refresh interviewee list
                interviewees = db_manager.get_personas_by_category(category_key, "interviewee")
                interviewee_choices = []
                for p in interviewees:
                    short_desc = ""
                    if p.background:
                        first_sentence = p.background.split('.')[0]
                        short_desc = (first_sentence[:100] + '...') if len(first_sentence) > 100 else first_sentence
                    interviewee_choices.append(f"{p.name} - {short_desc}")
                
                console.print(f"[green]Added {len(search_results)} new {category_name} profiles.[/green]")
            else:
                console.print(f"[yellow]No matching profiles found. Try a different search query.[/yellow]")
        
        selected_interviewee_displays = questionary.checkbox(
            f"Select interviewee personas (up to 10) for {category_name}:",
            choices=interviewee_choices
        ).ask()
        
        selected_interviewee_names = [display.split(' - ')[0] for display in selected_interviewee_displays]
        
        # Generate interviews
        console.print(f"\nGenerating interviews for {len(selected_interviewee_names)} personas in {category_name} category...")
        
        # Ask if user wants to use batch processing
        use_batch = questionary.confirm(
            "Do you want to use batch processing for faster generation (experimental)?",
            default=False
        ).ask()
        
        if use_batch:
            # Initialize batch processor
            batch_processor = BatchProcessor(model_info["provider"], model_info["model"])
            
            # Prepare prompts for all interviews
            interview_prompts = []
            for interviewee_name in selected_interviewee_names:
                interviewee = next(p for p in interviewees if p.name == interviewee_name)
                
                prompt = INTERVIEW_GENERATION_PROMPT.format(
                    interviewer_details=selected_interviewer.background,
                    interviewee_details=interviewee.background,
                    script=script
                )
                interview_prompts.append((interviewee, prompt))
            
            # Process interviews in batch
            console.print(f"Generating {len(interview_prompts)} interviews in batch...")
            prompts_only = [p[1] for p in interview_prompts]
            
            # Determine model's token limit from MODEL_INFO
            model_info_data = get_all_model_info().get(model_info["model"], {})
            token_limit = model_info_data.get('token_limit', 2048)  # Default to conservative 2048 if unknown
                    
            # Use 90% of the model's token limit for generation to be safe
            safe_token_limit = int(token_limit * 0.9)
            console.print(f"Using {safe_token_limit} tokens for generation (model max: {token_limit})")
                    
            results = batch_processor.process_batch_sync(prompts_only, max_tokens=safe_token_limit)
            
            # Process results
            for i, (interviewee, _) in enumerate(interview_prompts):
                interview_text = results[i]
                
                # Format as XML
                xml_prompt = XML_FORMATTING_PROMPT.format(
                    interview_id=f"{category_key}_{interviewee.id}",
                    interviewer_details=selected_interviewer.background,
                    interviewee_details=interviewee.background,
                    interview_text=interview_text
                )
                
                console.print(f"Formatting interview with {interviewee.name} as XML...")
                
                # Use a smaller token limit for XML formatting than for the main interview
                # This helps with models that have smaller context windows
                model_info_data = get_all_model_info().get(model_info["model"], {})
                xml_token_limit = min(model_info_data.get('token_limit', 4096), 4000)
                console.print(f"Using {xml_token_limit} tokens for XML formatting")
                
                xml_formatted = ai_model.generate_text(xml_prompt, max_tokens=xml_token_limit)
                
                # Save interview to database
                interview_id = db_manager.create_interview(
                    interviewer_id=selected_interviewer.id,
                    interviewee_id=interviewee.id,
                    category=category_key,
                    model_used=f"{model_info['provider']}/{model_info['model']}",
                    raw_interview=interview_text,
                    xml_formatted=xml_formatted
                )
                
                # Generate structured analysis with improved components
                # Create a more detailed analysis prompt that requests specific sections
                structured_analysis_prompt = f"""
Analyze the following interview between {selected_interviewer.name} and {interviewee.name} about AI in consulting.
Provide a structured analysis with the following sections:

INTERVIEW TEXT:
{interview_text}

Please format your analysis with these clear sections:

1. KEY POINTS: Summarize the 3-5 most important points from the interview.

2. NOTABLE QUOTES: Extract 2-3 direct quotes that best represent the interviewee's perspective.

3. AI ATTITUDES: Analyze the interviewee's attitude toward AI in consulting (positive, negative, neutral, nuanced).

4. RQ1 INSIGHTS: What insights does this interview provide about the state of AI adoption in consulting?

5. RQ2 INSIGHTS: What insights does this interview provide about current market trends in consulting?

6. RQ3 INSIGHTS: What insights does this interview provide about automation and knowledge management?

7. RQ4 INSIGHTS: What insights does this interview provide about ethical considerations and risks?

8. CONTRADICTIONS: Note any contradictions or inconsistencies in the interviewee's statements.

9. AUTHENTICITY ASSESSMENT: Evaluate how authentic and realistic this interview feels.
"""
                
                console.print(f"Generating structured analysis for {interviewee.name}...")
                
                # Use an appropriate token limit for analysis as well
                model_info_data = get_all_model_info().get(model_info["model"], {})
                analysis_token_limit = min(model_info_data.get('token_limit', 4096), 4000)
                console.print(f"Using {analysis_token_limit} tokens for analysis generation")
                
                analysis_text = ai_model.generate_text(structured_analysis_prompt, max_tokens=analysis_token_limit)
                
                # Parse the structured analysis into components
                import re
                
                # Create a parser to extract each section
                def extract_section(text, section_name):
                    pattern = rf"{section_name}:?(.*?)(?=\d+\.\s+\w+:|$)"
                    match = re.search(pattern, text, re.DOTALL | re.IGNORECASE)
                    if match:
                        return match.group(1).strip()
                    return ""
                
                # Extract each section
                analysis_data = {
                    "key_points": extract_section(analysis_text, r"1\.?\s*KEY POINTS"),
                    "notable_quotes": extract_section(analysis_text, r"2\.?\s*NOTABLE QUOTES"),
                    "ai_attitudes": extract_section(analysis_text, r"3\.?\s*AI ATTITUDES"),
                    "rq1_insights": extract_section(analysis_text, r"4\.?\s*RQ1 INSIGHTS"),
                    "rq2_insights": extract_section(analysis_text, r"5\.?\s*RQ2 INSIGHTS"),
                    "rq3_insights": extract_section(analysis_text, r"6\.?\s*RQ3 INSIGHTS"),
                    "rq4_insights": extract_section(analysis_text, r"7\.?\s*RQ4 INSIGHTS"),
                    "contradictions": extract_section(analysis_text, r"8\.?\s*CONTRADICTIONS"),
                    "authenticity_assessment": extract_section(analysis_text, r"9\.?\s*AUTHENTICITY")
                }
                
                # Fallback if structured parsing fails
                if not any(analysis_data.values()):
                    console.print("[yellow]Structured analysis parsing failed, using raw analysis[/yellow]")
                    analysis_data = {
                        "key_points": analysis_text,
                        "notable_quotes": "",
                        "ai_attitudes": "",
                        "rq1_insights": "",
                        "rq2_insights": "",
                        "rq3_insights": "",
                        "rq4_insights": "",
                        "contradictions": "",
                        "authenticity_assessment": ""
                    }
                
                db_manager.create_analysis(interview_id, analysis_data)
                
                # Increment counters
                category_interviews_generated += 1
                total_interviews_generated += 1
                
        else:
            # Process interviews sequentially 
            with Progress() as progress:
                tasks = {}
                for name in selected_interviewee_names:
                    task_id = progress.add_task(f"[green]Generating interview with {name}...", total=3)
                    tasks[name] = task_id
                
                for interviewee_name in selected_interviewee_names:
                    interviewee = next(p for p in interviewees if p.name == interviewee_name)
                    
                    # 1. Generate interview
                    prompt = INTERVIEW_GENERATION_PROMPT.format(
                        interviewer_details=selected_interviewer.background,
                        interviewee_details=interviewee.background,
                        script=script
                    )
                    
                    console.print(f"\nGenerating interview between {selected_interviewer.name} and {interviewee.name}...")
                    interview_text = ai_model.generate_text(prompt, max_tokens=3000)
                    progress.update(tasks[interviewee_name], advance=1)
                    
                    # 2. Format as XML
                    xml_prompt = XML_FORMATTING_PROMPT.format(
                        interview_id=f"{category_key}_{interviewee.id}",
                        interviewer_details=selected_interviewer.background,
                        interviewee_details=interviewee.background,
                        interview_text=interview_text
                    )
                    
                    console.print(f"Formatting interview as XML...")
                    xml_formatted = ai_model.generate_text(xml_prompt)
                    progress.update(tasks[interviewee_name], advance=1)
                    
                    # 3. Save interview to database
                    interview_id = db_manager.create_interview(
                        interviewer_id=selected_interviewer.id,
                        interviewee_id=interviewee.id,
                        category=category_key,
                        model_used=f"{model_info['provider']}/{model_info['model']}",
                        raw_interview=interview_text,
                        xml_formatted=xml_formatted
                    )
                    
                    # 4. Generate analysis
                    analysis_prompt = ANALYSIS_PROMPT.format(
                        interview_text=interview_text
                    )
                    
                    console.print(f"Generating analysis...")
                    analysis_text = ai_model.generate_text(analysis_prompt)
                    
                    # Parse analysis into components (simplified)
                    analysis_data = {
                        "key_points": analysis_text,
                        "notable_quotes": "",
                        "ai_attitudes": "",
                        "rq1_insights": "",
                        "rq2_insights": "",
                        "rq3_insights": "",
                        "rq4_insights": "",
                        "contradictions": "",
                        "authenticity_assessment": ""
                    }
                    
                    db_manager.create_analysis(interview_id, analysis_data)
                    
                    # Increment counters
                    total_interviews_generated += 1
                    
                    # Update progress
                    progress.update(tasks[interviewee_name], advance=1)
        
                # Mark this category as interviewed in this round
                interviewed_categories.add(category_key)
                
                console.print(f"[green]Generated {category_interviews_generated} interviews for {category_name}![/green]")
                
                # Ask if user wants to continue with next category or finish
                if category_name != selected_categories[-1]:  # Not the last category in this batch
                    continue_with_next = questionary.confirm(
                        f"Continue to next selected category?",
                        default=True
                    ).ask()
                    
                    if not continue_with_next:
                        break
        
    except Exception as e:
        console.print(f"[red]Error generating interviews: {str(e)}[/red]")
        
        # Ask if user wants to continue despite error
        continue_after_error = questionary.confirm(
            "Would you like to continue with the next categories despite the error?",
            default=True
        ).ask()
        
        if not continue_after_error:
            break
    
    # Final summary
    console.print("\n[bold green]Interview Generation Complete[/bold green]")
    console.print(f"Total interviews generated: {total_interviews_generated}")
    console.print(f"Categories covered: {len(interviewed_categories)}")
    
    # Prompt to export results if not already doing a report
    try:
        if total_interviews_generated > 0 and ('generate_report_now' not in locals() or not generate_report_now):
            export_now = questionary.confirm(
                "Would you like to export all interviews now?",
                default=True
            ).ask()
            
            if export_now:
                export_format = questionary.select(
                    "Select export format:",
                    choices=["all", "xml", "txt", "md", "json"]
                ).ask()
                
                include_analysis = questionary.confirm(
                    "Include analysis in exported files?",
                    default=True
                ).ask()
                
                # Call the export function
                export_interviews(format=export_format, output_dir='exports', include_analysis=include_analysis)
    except Exception as e:
        console.print(f"[red]Error exporting interviews: {str(e)}[/red]")

@cli.command()
@click.option('--format', type=click.Choice(['xml', 'txt', 'md', 'json', 'all']), default='all')
@click.option('--output-dir', default='exports')
@click.option('--include-analysis', is_flag=True, help="Include analysis in the exported files")
def export_interviews(format, output_dir, include_analysis):
    """Export generated interviews to files with timestamps and stakeholder information."""
    # Get all interviews from the database
    try:
        # Create output directory if it doesn't exist
        os.makedirs(output_dir, exist_ok=True)
        
        # Get current timestamp for folder organization
        import datetime
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        export_folder = os.path.join(output_dir, f"export_{timestamp}")
        os.makedirs(export_folder, exist_ok=True)
        
        # Create subdirectories by category
        for category in NORMALIZED_CATEGORIES.values():
            os.makedirs(os.path.join(export_folder, category), exist_ok=True)
        
        # Get interviews from database with all related info
        from sqlalchemy.orm import joinedload
        
        # This should be implemented to fetch all interviews
        # with related persona and analysis data
        interviews = db_manager.get_all_interviews_with_details()
        
        if not interviews:
            console.print("[yellow]No interviews found in the database.[/yellow]")
            return
        
        # Track stats for summary
        export_stats = {
            "total": len(interviews),
            "by_category": {},
            "by_model": {}
        }
        
        # Export each interview
        with Progress() as progress:
            task = progress.add_task("[green]Exporting interviews...", total=len(interviews))
            
            for interview in interviews:
                # Extract metadata for file naming
                interviewer = interview.interviewer
                interviewee = interview.interviewee
                category = interview.category
                model_used = interview.model_used.replace('/', '_')
                
                # Create a standardized base filename with rich metadata
                base_filename = f"{timestamp}_{category}_{interviewee.name.replace(' ', '_')}_{model_used}"
                
                # Update stats
                export_stats["by_category"][category] = export_stats["by_category"].get(category, 0) + 1
                export_stats["by_model"][model_used] = export_stats["by_model"].get(model_used, 0) + 1
                
                # Category-specific folder
                category_folder = os.path.join(export_folder, category)
                
                # Export in all requested formats
                formats_to_export = [format] if format != 'all' else ['xml', 'txt', 'md', 'json']
                
                for fmt in formats_to_export:
                    if fmt == 'xml':
                        # Export XML format - structured for research analysis
                        xml_path = os.path.join(category_folder, f"{base_filename}.xml")
                        with open(xml_path, 'w', encoding='utf-8') as f:
                            f.write(interview.xml_formatted)
                    
                    elif fmt == 'txt':
                        # Export plain text format - easy to read
                        txt_path = os.path.join(category_folder, f"{base_filename}.txt")
                        with open(txt_path, 'w', encoding='utf-8') as f:
                            f.write(f"INTERVIEW: {interviewee.name} ({category})\n")
                            f.write(f"INTERVIEWER: {interviewer.name}\n")
                            f.write(f"MODEL: {interview.model_used}\n")
                            f.write(f"DATE: {timestamp}\n\n")
                            f.write(interview.raw_interview)
                    
                    elif fmt == 'md':
                        # Export markdown format - good for documentation
                        md_path = os.path.join(category_folder, f"{base_filename}.md")
                        with open(md_path, 'w', encoding='utf-8') as f:
                            f.write(f"# Interview with {interviewee.name}\n\n")
                            f.write(f"**Category:** {category}\n")
                            f.write(f"**Interviewer:** {interviewer.name}\n")
                            f.write(f"**Model Used:** {interview.model_used}\n")
                            f.write(f"**Date:** {timestamp}\n\n")
                            f.write("## Interview Content\n\n")
                            
                            # Format the raw interview as markdown
                            lines = interview.raw_interview.split('\n')
                            for line in lines:
                                if line.startswith(f"{interviewer.name}:"):
                                    f.write(f"**{line}**\n\n")
                                elif line.startswith(f"{interviewee.name}:"):
                                    f.write(f"_{line}_\n\n")
                                else:
                                    f.write(f"{line}\n\n")
                            
                            # Include analysis if requested
                            if include_analysis and interview.analysis:
                                f.write("\n## Analysis\n\n")
                                f.write(f"### Key Points\n\n{interview.analysis.key_points}\n\n")
                                
                                if interview.analysis.notable_quotes:
                                    f.write(f"### Notable Quotes\n\n{interview.analysis.notable_quotes}\n\n")
                                    
                                if interview.analysis.ai_attitudes:
                                    f.write(f"### AI Attitudes\n\n{interview.analysis.ai_attitudes}\n\n")
                                    
                                for rq_num in range(1, 5):
                                    rq_field = f"rq{rq_num}_insights"
                                    if getattr(interview.analysis, rq_field):
                                        f.write(f"### RQ{rq_num} Insights\n\n{getattr(interview.analysis, rq_field)}\n\n")
                    
                    elif fmt == 'json':
                        # Export JSON format - machine readable
                        import json
                        json_path = os.path.join(category_folder, f"{base_filename}.json")
                        
                        # Construct interview data
                        interview_data = {
                            "metadata": {
                                "id": interview.id,
                                "timestamp": timestamp,
                                "category": category,
                                "model_used": interview.model_used
                            },
                            "interviewer": {
                                "id": interviewer.id,
                                "name": interviewer.name,
                                "background": interviewer.background
                            },
                            "interviewee": {
                                "id": interviewee.id,
                                "name": interviewee.name,
                                "category": interviewee.category,
                                "background": interviewee.background
                            },
                            "interview": {
                                "raw": interview.raw_interview,
                                "xml": interview.xml_formatted
                            }
                        }
                        
                        # Add analysis if requested
                        if include_analysis and interview.analysis:
                            interview_data["analysis"] = {
                                "key_points": interview.analysis.key_points,
                                "notable_quotes": interview.analysis.notable_quotes,
                                "ai_attitudes": interview.analysis.ai_attitudes,
                                "rq1_insights": interview.analysis.rq1_insights,
                                "rq2_insights": interview.analysis.rq2_insights,
                                "rq3_insights": interview.analysis.rq3_insights,
                                "rq4_insights": interview.analysis.rq4_insights,
                                "contradictions": interview.analysis.contradictions,
                                "authenticity_assessment": interview.analysis.authenticity_assessment
                            }
                        
                        with open(json_path, 'w', encoding='utf-8') as f:
                            json.dump(interview_data, f, indent=2)
                
                progress.update(task, advance=1)
        
        # Create a summary report
        summary_path = os.path.join(export_folder, "export_summary.md")
        with open(summary_path, 'w', encoding='utf-8') as f:
            f.write(f"# Interview Export Summary\n\n")
            f.write(f"**Date:** {timestamp}\n\n")
            f.write(f"**Total Interviews:** {export_stats['total']}\n\n")
            
            f.write("## Interviews by Category\n\n")
            for category, count in export_stats["by_category"].items():
                f.write(f"- {category}: {count}\n")
            
            f.write("\n## Interviews by Model\n\n")
            for model, count in export_stats["by_model"].items():
                f.write(f"- {model}: {count}\n")
            
            f.write("\n## Export Details\n\n")
            f.write(f"- Formats: {', '.join(formats_to_export)}\n")
            f.write(f"- Analysis Included: {include_analysis}\n")
            f.write(f"- Export Directory: {export_folder}\n")
        
        console.print(f"[green]Successfully exported {len(interviews)} interviews to {export_folder}/[/green]")
        console.print(f"[green]Export summary available at {summary_path}[/green]")
    
    except Exception as e:
        console.print(f"[red]Error exporting interviews: {str(e)}[/red]")

@cli.command()
@click.option('--output-format', type=click.Choice(['md', 'html', 'pdf', 'pptx', 'all']), default='md')
@click.option('--output-dir', default='exports')
@click.option('--individual-reports/--no-individual-reports', default=True, help="Generate detailed reports for each interview")
@click.option('--stakeholder-reports/--no-stakeholder-reports', default=True, help="Generate aggregated reports by stakeholder type")
@click.option('--summary-level', type=click.Choice(['basic', 'detailed', 'comprehensive']), default='detailed')
def generate_report(output_format, output_dir, individual_reports, stakeholder_reports, summary_level):
    """Generate multi-level reports from individual interviews to executive summaries."""
    try:
        # Get all interviews with analysis from the database
        interviews = db_manager.get_all_interviews_with_details()
        
        if not interviews or not any(hasattr(i, 'analysis') and i.analysis for i in interviews):
            console.print("[yellow]No interview analyses found in the database.[/yellow]")
            return
        
        console.print(f"[cyan]Found {len(interviews)} interviews to analyze[/cyan]")
        console.print("Generating multi-level analysis reports:")
        console.print(f"  • Individual reports: {'Yes' if individual_reports else 'No'}")
        console.print(f"  • Stakeholder group reports: {'Yes' if stakeholder_reports else 'No'}")
        console.print(f"  • Summary level: {summary_level}")
        console.print(f"  • Output format: {output_format}\n")
        
        # Create timestamp for the report
        import datetime
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        
        # Create report directories
        report_dir = os.path.join(output_dir, f'reports_{timestamp}')
        os.makedirs(report_dir, exist_ok=True)
        
        # Create subdirectories for different report types
        individual_dir = os.path.join(report_dir, 'individual')
        stakeholder_dir = os.path.join(report_dir, 'stakeholder_groups')
        summary_dir = os.path.join(report_dir, 'summary')
        presentation_dir = os.path.join(report_dir, 'presentation')
        
        os.makedirs(individual_dir, exist_ok=True)
        os.makedirs(stakeholder_dir, exist_ok=True)
        os.makedirs(summary_dir, exist_ok=True)
        os.makedirs(presentation_dir, exist_ok=True)
        
        # Group interviews by category
        interviews_by_category = {}
        
        # Also track by interviewee for individual reports
        interviews_by_interviewee = {}
        
        # Process all interviews
        for interview in interviews:
            if not hasattr(interview, 'analysis') or not interview.analysis:
                console.print(f"[yellow]Warning: Interview {interview.id} has no analysis and will be excluded[/yellow]")
                continue
                
            # Group by category
            category = interview.category
            if category not in interviews_by_category:
                interviews_by_category[category] = []
            interviews_by_category[category].append(interview)
            
            # Group by interviewee
            interviewee_id = interview.interviewee_id
            if interviewee_id not in interviews_by_interviewee:
                interviews_by_interviewee[interviewee_id] = []
            interviews_by_interviewee[interviewee_id].append(interview)
        
        console.print(f"Grouped interviews into {len(interviews_by_category)} stakeholder categories")
        
        # =========================================
        # 1. GENERATE INDIVIDUAL INTERVIEW REPORTS
        # =========================================
        if individual_reports:
            console.print("\n[bold cyan]Generating individual interview reports...[/bold cyan]")
            
            with Progress() as progress:
                task = progress.add_task("[green]Processing individual interviews...", total=len(interviews))
                
                for interview in interviews:
                    # Skip interviews without analysis
                    if not hasattr(interview, 'analysis') or not interview.analysis:
                        progress.update(task, advance=1)
                        continue
                    
                    # Generate a detailed report for this individual interview
                    report_content = f"# Interview Analysis: {interview.interviewee.name} ({interview.category})\n\n"
                    report_content += f"**Date:** {interview.created_at.strftime('%Y-%m-%d')}\n"
                    report_content += f"**Interviewer:** {interview.interviewer.name}\n"
                    report_content += f"**Model Used:** {interview.model_used}\n\n"
                    
                    # Include the key analysis sections
                    if interview.analysis.key_points:
                        report_content += f"## Key Points\n\n{interview.analysis.key_points}\n\n"
                    
                    if interview.analysis.notable_quotes:
                        report_content += f"## Notable Quotes\n\n{interview.analysis.notable_quotes}\n\n"
                    
                    if interview.analysis.ai_attitudes:
                        report_content += f"## Attitudes Toward AI\n\n{interview.analysis.ai_attitudes}\n\n"
                    
                    # Include research question insights
                    report_content += "## Research Question Insights\n\n"
                    
                    if interview.analysis.rq1_insights:
                        report_content += f"### RQ1: AI Adoption\n\n{interview.analysis.rq1_insights}\n\n"
                    
                    if interview.analysis.rq2_insights:
                        report_content += f"### RQ2: Market Trends\n\n{interview.analysis.rq2_insights}\n\n"
                    
                    if interview.analysis.rq3_insights:
                        report_content += f"### RQ3: Automation & Knowledge\n\n{interview.analysis.rq3_insights}\n\n"
                    
                    if interview.analysis.rq4_insights:
                        report_content += f"### RQ4: Ethical Considerations\n\n{interview.analysis.rq4_insights}\n\n"
                    
                    # Add additional analysis if available
                    if interview.analysis.contradictions:
                        report_content += f"## Contradictions & Inconsistencies\n\n{interview.analysis.contradictions}\n\n"
                    
                    if interview.analysis.authenticity_assessment:
                        report_content += f"## Authenticity Assessment\n\n{interview.analysis.authenticity_assessment}\n\n"
                    
                    # Include a sample of the raw interview
                    report_content += "## Interview Excerpt\n\n"
                    
                    # Include first 1000 characters of the interview as a sample
                    excerpt = interview.raw_interview[:1000] + "..." if len(interview.raw_interview) > 1000 else interview.raw_interview
                    report_content += f"```\n{excerpt}\n```\n\n"
                    
                    # Save the individual report
                    interviewee_name = interview.interviewee.name.replace(' ', '_')
                    report_filename = f"{interview.category}_{interviewee_name}_{interview.id}.md"
                    with open(os.path.join(individual_dir, report_filename), 'w', encoding='utf-8') as f:
                        f.write(report_content)
                    
                    progress.update(task, advance=1)
            
            console.print(f"[green]Generated {len(interviews)} individual interview reports in {individual_dir}[/green]")
        
        # =============================================
        # 2. GENERATE STAKEHOLDER GROUP ANALYSIS REPORTS
        # =============================================
        if stakeholder_reports:
            console.print("\n[bold cyan]Generating stakeholder group reports...[/bold cyan]")
            
            # Prepare AI model for synthesis
            model_info = select_ai_model()
            ai_model = AIModelInterface(model_info["provider"], model_info["model"])
            
            with Progress() as progress:
                task = progress.add_task("[green]Processing stakeholder groups...", total=len(interviews_by_category))
                
                for category, category_interviews in interviews_by_category.items():
                    # Generate a comprehensive analysis for this stakeholder group
                    report_content = f"# {category.replace('_', ' ').title()} Analysis Report\n\n"
                    report_content += f"**Generated:** {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
                    report_content += f"## Overview\n\n"
                    report_content += f"This report analyzes {len(category_interviews)} interviews with {category.replace('_', ' ')} stakeholders.\n\n"
                    
                    # Add table of interviewees
                    report_content += "### Interviews Included\n\n"
                    report_content += "| Interviewee | Model Used | Date |\n"
                    report_content += "|-------------|-----------|------|\n"
                    
                    for interview in category_interviews:
                        interview_date = interview.created_at.strftime('%Y-%m-%d')
                        report_content += f"| {interview.interviewee.name} | {interview.model_used} | {interview_date} |\n"
                    
                    report_content += "\n"
                    
                    # Collect all insights by research question
                    rq_insights = {
                        "rq1": [i.analysis.rq1_insights for i in category_interviews if hasattr(i, 'analysis') and i.analysis and i.analysis.rq1_insights],
                        "rq2": [i.analysis.rq2_insights for i in category_interviews if hasattr(i, 'analysis') and i.analysis and i.analysis.rq2_insights],
                        "rq3": [i.analysis.rq3_insights for i in category_interviews if hasattr(i, 'analysis') and i.analysis and i.analysis.rq3_insights],
                        "rq4": [i.analysis.rq4_insights for i in category_interviews if hasattr(i, 'analysis') and i.analysis and i.analysis.rq4_insights]
                    }
                    
                    # Collect all AI attitudes
                    attitudes = [i.analysis.ai_attitudes for i in category_interviews if hasattr(i, 'analysis') and i.analysis and i.analysis.ai_attitudes]
                    
                    # Extract all analysis objects for this category
                    category_analyses = [i.analysis for i in category_interviews 
                                        if hasattr(i, 'analysis') and i.analysis]
                    
                    # Generate executive summary and presentation bullets for this category
                    if category_analyses:
                        console.print(f"[cyan]Generating executive summary for {category} stakeholders...[/cyan]")
                        summary_results = generate_executive_summary(ai_model, category_analyses, category)
                        
                        # Add executive summary
                        report_content += "## Executive Summary\n\n"
                        report_content += summary_results["executive_summary"] + "\n\n"
                        
                        # Add presentation-ready bullets
                        report_content += "## Key Findings for Presentation\n\n"
                        report_content += "*The following bullet points are designed for direct use in presentation slides:*\n\n"
                        report_content += summary_results["presentation_bullets"] + "\n\n"
                    else:
                        # If no analyses, add placeholder section
                        report_content += "## Key Findings\n\n"
                        report_content += "No analysis data available for this stakeholder category.\n\n"
                    
                    # Generate synthesis for each research question
                    report_content += "## Research Question Insights\n\n"
                    
                    for rq_num, rq_name, rq_description in [
                        ("1", "AI Adoption", "How established is AI adoption within the consulting industry?"),
                        ("2", "Market Trends", "What are the current market trends in the consulting industry?"),
                        ("3", "Automation & Knowledge", "How is AI affecting consulting firms in terms of automation and knowledge?"),
                        ("4", "Ethical Considerations", "What ethical risks are associated with AI in consulting?")
                    ]:
                        report_content += f"### RQ{rq_num}: {rq_name}\n\n"
                        report_content += f"*{rq_description}*\n\n"
                        
                        insights = rq_insights[f"rq{rq_num}"]
                        
                        if insights and len(insights) > 1:
                            # Create a synthesis prompt for this research question
                            rq_prompt = f"""
Synthesize these insights about {rq_name} from {len(insights)} {category.replace('_', ' ')} stakeholders:

{chr(10).join(insights[:3])}

Create a coherent summary that:
1. Identifies the main themes
2. Organizes the insights into logical categories
3. Highlights significant agreements and disagreements
4. Draws meaningful conclusions

Be concise but comprehensive (250-300 words).
"""
                            # Generate the RQ synthesis
                            rq_synthesis = ai_model.generate_text(rq_prompt)
                            report_content += rq_synthesis + "\n\n"
                        elif insights:
                            # For a single insight, use it directly
                            report_content += insights[0] + "\n\n"
                        else:
                            report_content += "No specific insights available for this research question.\n\n"
                    
                    # Add notable quotes section
                    report_content += "## Notable Quotes\n\n"
                    
                    for interview in category_interviews:
                        if hasattr(interview, 'analysis') and interview.analysis and interview.analysis.notable_quotes:
                            report_content += f"**{interview.interviewee.name}:**\n\n"
                            report_content += f"{interview.analysis.notable_quotes}\n\n"
                    
                    # Add AI attitudes synthesis
                    if attitudes and len(attitudes) > 1:
                        report_content += "## Attitudes Toward AI\n\n"
                        
                        # Create a synthesis prompt for AI attitudes
                        attitudes_prompt = f"""
Analyze these attitudes toward AI from {len(attitudes)} {category.replace('_', ' ')} stakeholders:

{chr(10).join(attitudes[:3])}

Create a summary that:
1. Identifies the range of attitudes (positive, negative, neutral, nuanced)
2. Explains the reasons behind these attitudes
3. Categorizes common concerns and expectations
4. Draws conclusions about how this stakeholder group views AI

Keep your response to 200-250 words.
"""
                        # Generate the attitudes synthesis
                        attitudes_synthesis = ai_model.generate_text(attitudes_prompt)
                        report_content += attitudes_synthesis + "\n\n"
                    elif attitudes:
                        report_content += "## Attitudes Toward AI\n\n"
                        report_content += attitudes[0] + "\n\n"
                    
                    # Save the stakeholder group report in all requested formats
                    convert_report_to_formats(
                        report_content=report_content,
                        base_filename=f"{category}_analysis",
                        output_dir=stakeholder_dir,
                        formats=output_format
                    )
                    
                    # Also create a PowerPoint for this stakeholder group if we have presentation bullets
                    if category_analyses and "presentation_bullets" in summary_results:
                        # Create a category-specific presentation directory
                        category_presentation_dir = os.path.join(presentation_dir, category)
                        os.makedirs(category_presentation_dir, exist_ok=True)
                        
                        # Create the PowerPoint presentation
                        create_presentation(
                            presentation_bullets=summary_results["presentation_bullets"],
                            output_dir=category_presentation_dir,
                            title=f"{category.replace('_', ' ').title()} Stakeholder Analysis",
                            timestamp=timestamp,
                            filename=f"{category}_analysis"
                        )
                    
                    progress.update(task, advance=1)
            
            console.print(f"[green]Generated {len(interviews_by_category)} stakeholder group reports in {stakeholder_dir}[/green]")
        
        # =============================================
        # 3. GENERATE COMPREHENSIVE SUMMARY REPORT
        # =============================================
        
        # Start building the comprehensive report
        report_content = f"# AI in Consulting Research Report\n\n"
        report_content += f"**Generated:** {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
        report_content += f"## Overview\n\n"
        report_content += f"This report analyzes {len(interviews)} interviews across {len(interviews_by_category)} stakeholder categories.\n\n"
        
        # Add table of contents
        report_content += "## Table of Contents\n\n"
        report_content += "1. [Executive Summary](#executive-summary)\n"
        report_content += "2. [Key Findings for Presentation](#key-findings-for-presentation)\n"
        report_content += "3. [Research Questions](#research-questions)\n"
        for i, category in enumerate(interviews_by_category.keys()):
            report_content += f"{i+4}. [Analysis: {category.replace('_', ' ').title()}](#{category})\n"
        report_content += f"{len(interviews_by_category)+4}. [Cross-Category Insights](#cross-category-insights)\n"
        report_content += f"{len(interviews_by_category)+5}. [Methodology](#methodology)\n\n"
        
        # Generate executive summary using AI
        console.print("[cyan]Generating executive summary...[/cyan]")
        
        # Prepare AI model for report generation
        model_info = select_ai_model()
        ai_model = AIModelInterface(model_info["provider"], model_info["model"])
        
        # Extract all analysis objects
        all_analyses = [i.analysis for i in interviews if hasattr(i, 'analysis') and i.analysis]
        
        # Generate comprehensive executive summary and presentation bullets
        if all_analyses:
            # Generate the executive summary and presentation-ready bullets
            summary_results = generate_executive_summary(ai_model, all_analyses)
            
            # Add executive summary
            report_content += "## Executive Summary\n\n"
            report_content += summary_results["executive_summary"] + "\n\n"
            
            # Add presentation-ready bullets
            report_content += "## Key Findings for Presentation<a name='key-findings-for-presentation'></a>\n\n"
            report_content += "*The following bullet points are designed for direct use in presentation slides:*\n\n"
            report_content += summary_results["presentation_bullets"] + "\n\n"
            
                # Create PowerPoint presentation using the helper function
            create_presentation(
                presentation_bullets=summary_results["presentation_bullets"],
                output_dir=presentation_dir,
                title="AI in Consulting Research Findings",
                timestamp=timestamp,
                filename=f"key_findings_{timestamp}"
            )
        else:
            report_content += "## Executive Summary\n\n"
            report_content += "No analyses available to generate an executive summary.\n\n"
            report_content += "## Key Findings for Presentation\n\n"
            report_content += "No analyses available to generate presentation points.\n\n"
        
        # Add research questions section
        report_content += "## Research Questions\n\n"
        report_content += "This research examines the following key questions:\n\n"
        report_content += "1. **RQ1:** How established is AI adoption within the consulting industry?\n\n"
        report_content += "2. **RQ2:** What are the current market trends in the consulting industry, particularly in relation to AI?\n\n"
        report_content += "3. **RQ3:** How is AI affecting consulting firms in terms of automation and knowledge internalization?\n\n"
        report_content += "4. **RQ4:** What ethical risks are associated with the integration of AI in consulting practices?\n\n"
        
        # Generate category analyses
        for category, category_interviews in interviews_by_category.items():
            report_content += f"## Analysis: {category.replace('_', ' ').title()}<a name='{category}'></a>\n\n"
            
            # Add category overview
            report_content += f"### Overview\n\n"
            report_content += f"This section analyzes {len(category_interviews)} interviews with {category.replace('_', ' ')}.\n\n"
            
            # Aggregate insights for each research question
            rq_insights = {
                "rq1": [],
                "rq2": [],
                "rq3": [],
                "rq4": []
            }
            
            # Collect attitudes toward AI
            attitudes = []
            
            for interview in category_interviews:
                if hasattr(interview, 'analysis') and interview.analysis:
                    analysis = interview.analysis
                    
                    if analysis.rq1_insights:
                        rq_insights["rq1"].append(analysis.rq1_insights)
                    if analysis.rq2_insights:
                        rq_insights["rq2"].append(analysis.rq2_insights)
                    if analysis.rq3_insights:
                        rq_insights["rq3"].append(analysis.rq3_insights)
                    if analysis.rq4_insights:
                        rq_insights["rq4"].append(analysis.rq4_insights)
                    if analysis.ai_attitudes:
                        attitudes.append(analysis.ai_attitudes)
            
            # Generate a consolidated analysis for each research question
            for rq in ["rq1", "rq2", "rq3", "rq4"]:
                if rq_insights[rq]:
                    # Generate a synthesis of insights for this RQ using AI
                    synthesis_prompt = f"""
Synthesize the following insights about {rq.upper()} (Research Question {rq[2]}) 
from multiple {category.replace('_', ' ')} interviews:

INSIGHTS:
{"".join(rq_insights[rq][:3])}  # Use first 3 to avoid token limits

Create a concise summary (250 words) that:
1. Identifies common themes and patterns
2. Highlights unique perspectives
3. Draws meaningful conclusions from these insights
"""
                    
                    synthesis = ai_model.generate_text(synthesis_prompt)
                    
                    # Add to report
                    report_content += f"### {rq.upper()} Synthesis\n\n"
                    report_content += f"{synthesis}\n\n"
            
            # Add attitudes toward AI section if available
            if attitudes:
                attitude_prompt = f"""
Summarize the attitudes toward AI in consulting based on the following assessments
from {category.replace('_', ' ')} interviews:

ATTITUDES:
{"".join(attitudes[:3])}  # Use first 3 to avoid token limits

Create a concise summary (150 words) that categorizes and explains the range of attitudes observed.
"""
                
                attitude_synthesis = ai_model.generate_text(attitude_prompt)
                report_content += f"### Attitudes Toward AI\n\n"
                report_content += f"{attitude_synthesis}\n\n"
            
            # Add notable quotes section
            report_content += f"### Notable Quotes\n\n"
            quotes_added = False
            
            for interview in category_interviews:
                if hasattr(interview, 'analysis') and interview.analysis and interview.analysis.notable_quotes:
                    interviewee_name = interview.interviewee.name
                    quotes = interview.analysis.notable_quotes
                    report_content += f"**{interviewee_name}:**\n\n"
                    report_content += f"{quotes}\n\n"
                    quotes_added = True
            
            if not quotes_added:
                report_content += "No notable quotes available for this category.\n\n"
        
        # Generate cross-category analysis
        console.print("[cyan]Generating cross-category analysis...[/cyan]")
        
        # Prepare data for cross-category analysis prompt
        category_summaries = []
        for category in interviews_by_category.keys():
            category_summaries.append(f"- {category.replace('_', ' ').title()}")
        
        cross_category_prompt = f"""
Generate a cross-category analysis for a research report on "The Role of Business Consulting Firms in the Era of Artificial Intelligence"
comparing perspectives across these stakeholder groups:

{chr(10).join(category_summaries)}

Your analysis should:
1. Compare and contrast perspectives across stakeholder groups
2. Identify areas of consensus and disagreement
3. Highlight how different stakeholders view AI's impact on consulting
4. Draw connections between stakeholder interests and attitudes
5. Be approximately 500 words in length
"""
        
        cross_category_analysis = ai_model.generate_text(cross_category_prompt)
        
        report_content += "## Cross-Category Insights<a name='cross-category-insights'></a>\n\n"
        report_content += cross_category_analysis + "\n\n"
        
        # Add methodology section
        report_content += "## Methodology<a name='methodology'></a>\n\n"
        report_content += "This research employed a virtual interview methodology, generating simulated conversations "
        report_content += "between AI-powered personas representing various stakeholders in the consulting ecosystem. "
        report_content += "Large language models were used to create realistic, diverse perspectives on the research questions.\n\n"
        
        report_content += "### Interview Distribution\n\n"
        report_content += "| Stakeholder Category | Interviews |\n"
        report_content += "|----------------------|------------|\n"
        
        for category, category_interviews in interviews_by_category.items():
            report_content += f"| {category.replace('_', ' ').title()} | {len(category_interviews)} |\n"
        
        report_content += "\n### Analysis Process\n\n"
        report_content += "Each interview underwent systematic analysis to extract key insights, notable quotes, and perspectives on the research questions. "
        report_content += "Cross-interview synthesis was performed to identify patterns and divergences across stakeholder groups.\n\n"
        
        # Add report footer
        report_content += "---\n\n"
        report_content += "*This report was generated using AI-powered analysis tools. "
        report_content += "While efforts have been made to ensure accuracy, the synthesized insights should be "
        report_content += "considered as exploratory findings rather than definitive conclusions.*\n"
        
        # Save the report
        report_path = os.path.join(report_dir, f'research_report_{timestamp}.md')
        with open(report_path, 'w', encoding='utf-8') as f:
            f.write(report_content)
        
        console.print(f"[green]Successfully generated report at {report_path}[/green]")
        
        # Convert to multiple formats
        convert_report_to_formats(
            report_content=report_content,
            base_filename=f'research_report_{timestamp}',
            output_dir=report_dir,
            formats=output_format
        )
        
def convert_report_to_formats(report_content, base_filename, output_dir, formats):
    """
    Convert a markdown report to various output formats.
    
    Args:
        report_content: The markdown content to convert
        base_filename: Base filename without extension
        output_dir: Directory to save the output files
        formats: Format or list of formats to convert to ('md', 'html', 'pdf', 'all')
    """
    # Save markdown version
    md_path = os.path.join(output_dir, f'{base_filename}.md')
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write(report_content)
    console.print(f"[green]Markdown report saved to {md_path}[/green]")
    
    # Determine which formats to generate
    formats_to_generate = []
    if formats == 'all':
        formats_to_generate = ['html', 'pdf']
    elif isinstance(formats, list):
        formats_to_generate = formats
    else:
        formats_to_generate = [formats]
        
    # Remove 'all' and 'md' from formats list since we already saved markdown
    formats_to_generate = [fmt for fmt in formats_to_generate if fmt not in ('all', 'md')]
    
    # Convert to HTML if requested
    if 'html' in formats_to_generate:
        try:
            import markdown
            html_path = os.path.join(output_dir, f'{base_filename}.html')
            with open(html_path, 'w', encoding='utf-8') as f:
                html_content = f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>AI in Consulting Research Report</title>
    <style>
        body {{
            font-family: Arial, sans-serif;
            line-height: 1.6;
            max-width: 900px;
            margin: 0 auto;
            padding: 20px;
        }}
        h1, h2, h3 {{
            color: #333;
        }}
        h1 {{
            border-bottom: 2px solid #333;
            padding-bottom: 10px;
        }}
        h2 {{
            margin-top: 30px;
            border-bottom: 1px solid #ccc;
            padding-bottom: 5px;
        }}
        blockquote {{
            border-left: 4px solid #ccc;
            padding-left: 15px;
            margin-left: 0;
            color: #555;
        }}
        table {{
            border-collapse: collapse;
            width: 100%;
            margin: 20px 0;
        }}
        th, td {{
            border: 1px solid #ddd;
            padding: 8px;
            text-align: left;
        }}
        th {{
            background-color: #f2f2f2;
        }}
        tr:nth-child(even) {{
            background-color: #f9f9f9;
        }}
    </style>
</head>
<body>
    {markdown.markdown(report_content, extensions=['tables'])}
</body>
</html>
"""
                f.write(html_content)
            console.print(f"[green]HTML version generated at {html_path}[/green]")
        except ImportError:
            console.print("[yellow]Could not generate HTML version. Install the 'markdown' package if needed.[/yellow]")
            console.print("Run: pip install markdown")
    
    # Convert to PDF if requested
    if 'pdf' in formats_to_generate:
        try:
            # Try using weasyprint for PDF conversion
            from weasyprint import HTML, CSS
            
            # First ensure we have HTML content
            import markdown
            html_content = f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>AI in Consulting Research Report</title>
    <style>
        @page {{
            margin: 1cm;
        }}
        body {{
            font-family: Arial, sans-serif;
            line-height: 1.6;
            margin: 1cm;
        }}
        h1, h2, h3 {{
            color: #333;
        }}
        h1 {{
            border-bottom: 2px solid #333;
            padding-bottom: 10px;
        }}
        h2 {{
            margin-top: 30px;
            border-bottom: 1px solid #ccc;
            padding-bottom: 5px;
        }}
        blockquote {{
            border-left: 4px solid #ccc;
            padding-left: 15px;
            margin-left: 0;
            color: #555;
        }}
        table {{
            border-collapse: collapse;
            width: 100%;
            margin: 20px 0;
        }}
        th, td {{
            border: 1px solid #ddd;
            padding: 8px;
            text-align: left;
        }}
        th {{
            background-color: #f2f2f2;
        }}
        tr:nth-child(even) {{
            background-color: #f9f9f9;
        }}
    </style>
</head>
<body>
    {markdown.markdown(report_content, extensions=['tables'])}
</body>
</html>
"""
            
            # Generate PDF
            pdf_path = os.path.join(output_dir, f'{base_filename}.pdf')
            HTML(string=html_content).write_pdf(pdf_path)
            console.print(f"[green]PDF version generated at {pdf_path}[/green]")
            
        except ImportError:
            # Alternative method using pdfkit if weasyprint is not available
            try:
                import pdfkit
                import markdown
                
                # Convert markdown to HTML
                html_content = markdown.markdown(report_content, extensions=['tables'])
                
                # Add styling
                styled_html = f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>AI in Consulting Research Report</title>
    <style>
        body {{
            font-family: Arial, sans-serif;
            line-height: 1.6;
            margin: 1cm;
        }}
        h1, h2, h3 {{
            color: #333;
        }}
        table {{
            border-collapse: collapse;
            width: 100%;
        }}
        th, td {{
            border: 1px solid #ddd;
            padding: 8px;
        }}
    </style>
</head>
<body>
    {html_content}
</body>
</html>
"""
                
                # Generate PDF
                pdf_path = os.path.join(output_dir, f'{base_filename}.pdf')
                pdfkit.from_string(styled_html, pdf_path)
                console.print(f"[green]PDF version generated at {pdf_path}[/green]")
                
            except ImportError:
                console.print("[yellow]Could not generate PDF version. Install 'weasyprint' or 'pdfkit' packages if needed.[/yellow]")
                console.print("Run: pip install weasyprint")
                console.print("or: pip install pdfkit")
            except Exception as e:
                console.print(f"[yellow]Error generating PDF: {str(e)}[/yellow]")
    
    except Exception as e:
        console.print(f"[red]Error generating report: {str(e)}[/red]")

if __name__ == '__main__':
    # Ensure directories exist
    os.makedirs('data/scripts', exist_ok=True)
    os.makedirs('data/personas', exist_ok=True)
    os.makedirs('data/interviews', exist_ok=True)
    os.makedirs('exports', exist_ok=True)
    
    # Show a welcome message
    console.print("\n[bold green]Virtual Interview Generator for AI in Consulting Research[/bold green]")
    console.print("This application helps you generate simulated interviews with various stakeholders")
    console.print("about the role of AI in the consulting industry. Follow the prompts to create")
    console.print("personas, generate interviews, and analyze the results.\n")
    
    cli()